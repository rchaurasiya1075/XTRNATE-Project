"""Excel → branded animated PPT. Any workbook, auto charts + tables."""
from __future__ import annotations

from datetime import datetime
from io import BytesIO
from zoneinfo import ZoneInfo

import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from utils.anim_deck import (
    GREEN,
    GREEN2,
    GOLD,
    MUTED,
    WHITE,
    INK,
    _add_chart,
    _add_gif,
    _appear_shapes,
    _fade_transition,
    _gif_bars,
    _gif_line,
    _header,
    _rect,
    _series,
    _tb,
)

IST = ZoneInfo("Asia/Kolkata")


def _cell(v) -> str:
    if v is None or (isinstance(v, float) and pd.isna(v)):
        return ""
    try:
        if pd.isna(v):
            return ""
    except Exception:
        pass
    if isinstance(v, datetime):
        return v.strftime("%d-%b-%Y")
    if isinstance(v, pd.Timestamp):
        if pd.isna(v):
            return ""
        return v.strftime("%d-%b-%Y")
    if isinstance(v, float):
        if v == int(v):
            return f"{int(v):,}"
        return f"{v:,.2f}"
    s = str(v).strip()
    return "" if s.lower() in ("nan", "none", "<na>") else s[:48]


def load_workbook(uploaded) -> dict[str, pd.DataFrame]:
    name = getattr(uploaded, "name", "book.xlsx").lower()
    if name.endswith(".csv"):
        df = pd.read_csv(uploaded)
        return {"Data": _clean(df)}
    xl = pd.ExcelFile(uploaded)
    out = {}
    for sn in xl.sheet_names[:8]:
        try:
            df = pd.read_excel(xl, sheet_name=sn)
        except Exception:
            continue
        df = _clean(df)
        if df is None or df.empty:
            continue
        out[str(sn)[:28]] = df
    return out


def _clean(df: pd.DataFrame) -> pd.DataFrame:
    if df is None or df.empty:
        return pd.DataFrame()
    df = df.copy()
    df.columns = [str(c).strip() if str(c).lower() != "nan" else f"Col{i+1}" for i, c in enumerate(df.columns)]
    df = df.dropna(how="all").dropna(axis=1, how="all")
    df = df.loc[:, ~df.columns.duplicated()]
    return df.reset_index(drop=True)


def _is_kpi_sheet(df: pd.DataFrame) -> bool:
    if df is None or df.empty or len(df.columns) < 2:
        return False
    if len(df) <= 12 and len(df.columns) <= 3:
        num = pd.to_numeric(df.iloc[:, 1], errors="coerce")
        return num.notna().mean() > 0.5
    return False


def _label_num(df: pd.DataFrame):
    if df is None or df.empty:
        return [], []
    return _series(df, df.columns[0], "Count" if "Count" in df.columns else df.columns[1] if len(df.columns) > 1 else df.columns[0], n=12)


def _add_table(slide, df, l, t, w, h, max_rows=11, max_cols=7):
    cols = list(df.columns)[:max_cols]
    rows = [cols]
    for _, r in df.head(max_rows).iterrows():
        rows.append([_cell(r[c]) for c in cols])
    nr, nc = len(rows), max(len(cols), 1)
    table = slide.shapes.add_table(nr, nc, Inches(l), Inches(t), Inches(w), Inches(h)).table
    for j in range(nc):
        table.columns[j].width = Inches(w / nc)
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = table.cell(i, j)
            cell.text = ""
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER if j else PP_ALIGN.LEFT
            run = p.add_run()
            run.text = str(val)
            run.font.size = Pt(10 if i else 11)
            run.font.bold = i == 0 or str(row[0]).upper() in ("TOTAL", "GRAND TOTAL", "TOTAL / AVERAGE")
            run.font.name = "Calibri"
            head = i == 0
            tot = str(row[0]).upper() in ("TOTAL", "GRAND TOTAL", "TOTAL / AVERAGE")
            run.font.color.rgb = WHITE if head or tot else INK
            cell.fill.solid()
            if head or tot:
                cell.fill.fore_color.rgb = GREEN
            elif i % 2:
                cell.fill.fore_color.rgb = RGBColor(0xF3, 0xF5, 0xF4)
            else:
                cell.fill.fore_color.rgb = WHITE


def workbook_to_pptx(sheets: dict, *, title="XTRNATE Report", subtitle="") -> bytes:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    stamp = datetime.now(IST).strftime("%d-%b-%Y %I:%M %p IST")
    sub = subtitle or stamp

    s = prs.slides.add_slide(blank)
    _rect(s, 0, 0, 13.333, 7.5, GREEN2)
    _rect(s, 0, 0, 0.16, 7.5, GOLD)
    _tb(s, 0.7, 1.6, 12, 0.35, "XTRNATE  •  EXCEL → PPT AUTOMATION", 14, True, GOLD)
    _tb(s, 0.7, 2.15, 12, 1.1, title[:48], 32, True, WHITE)
    _tb(s, 0.7, 3.4, 12, 0.4, sub, 16, False, MUTED)
    _tb(s, 0.7, 4.0, 12, 0.4, f"{len(sheets)} sheet(s)  •  Slideshow F5 for animated graphs", 14, False, WHITE)
    _tb(s, 0.7, 6.5, 12, 0.3, "Confidential", 12, False, MUTED)
    _fade_transition(s)

    # Index
    s = prs.slides.add_slide(blank)
    _header(s, "Workbook map")
    idx = pd.DataFrame(
        [{"Sheet": k, "Rows": len(v), "Columns": len(v.columns)} for k, v in sheets.items()]
    )
    _add_table(s, idx, 0.4, 1.15, 12.5, min(5.8, 0.42 * (len(idx) + 1) + 0.4))
    _fade_transition(s)

    for name, df in sheets.items():
        if df is None or df.empty:
            continue

        # KPI cards
        if _is_kpi_sheet(df):
            s = prs.slides.add_slide(blank)
            _header(s, f"{name}  •  snapshot")
            n = min(len(df), 5)
            for i in range(n):
                lab = _cell(df.iloc[i, 0])
                val = _cell(df.iloc[i, 1])
                left = 0.4 + i * 2.56
                _rect(s, left, 1.5, 2.4, 1.7, GREEN if i % 2 == 0 else RGBColor(0x2D, 0x6A, 0x4F))
                _tb(s, left + 0.08, 1.62, 2.24, 0.4, lab[:22], 11, True, MUTED, PP_ALIGN.CENTER)
                _tb(s, left + 0.08, 2.1, 2.24, 0.8, val[:16], 22, True, WHITE, PP_ALIGN.CENTER)
            _fade_transition(s)
            _appear_shapes(s, skip=1)

        labels, vals = _label_num(df)
        looks_date = False
        if labels:
            looks_date = sum(ch.isdigit() for ch in "".join(labels[:3])) >= 4 or any(
                x in str(df.columns[0]).lower() for x in ("date", "day", "month")
            )
        if labels and vals and max(vals) > 0:
            s = prs.slides.add_slide(blank)
            _header(s, f"{name}  •  animated chart")
            gif = _gif_line(labels, vals, name) if looks_date or len(labels) > 10 else _gif_bars(labels, vals, name)
            if gif:
                _add_gif(s, gif, 0.35, 1.05, 12.6, 5.9)
            else:
                _add_chart(s, "line" if looks_date else "bar", labels[:14], vals[:14], name[:20], 0.4, 1.15, 12.5, 5.7)
            _fade_transition(s)

        # Table
        s = prs.slides.add_slide(blank)
        _header(s, f"{name}  •  data")
        _add_table(s, df, 0.3, 1.05, 12.7, 5.9)
        _tb(s, 0.4, 7.05, 12, 0.25, f"Showing first rows  •  {len(df)} total", 10, False, INK)
        _fade_transition(s)

    s = prs.slides.add_slide(blank)
    _rect(s, 0, 0, 13.333, 7.5, GREEN2)
    _tb(s, 0.7, 2.5, 12, 0.5, "Excel → PPT complete", 28, True, WHITE)
    _tb(s, 0.7, 3.3, 12, 0.6, "F5 Slideshow  •  graphs animate  •  tables match forest-green format", 16, False, MUTED)
    _fade_transition(s)

    out = BytesIO()
    prs.save(out)
    return out.getvalue()
