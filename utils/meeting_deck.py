"""15-slide ISP review PPT for meetings."""
from io import BytesIO
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

NAVY = RGBColor(0x0B, 0x1F, 0x3A)
BLUE = RGBColor(0x0F, 0x4C, 0x81)
GOLD = RGBColor(0xD4, 0xA3, 0x37)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SLATE = RGBColor(0x1E, 0x29, 0x3B)
TEAL = RGBColor(0x0D, 0x94, 0x88)
LIGHT = RGBColor(0xF1, 0xF5, 0xF9)
MUTED = RGBColor(0x94, 0xA3, 0xB8)
RED = RGBColor(0xB9, 0x1C, 0x1C)


def _run(p, text, size=14, bold=False, color=WHITE):
    r = p.add_run()
    r.text = str(text)
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = "Calibri"
    return r


def _rect(slide, l, t, w, h, fill):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.fill.background()
    return sh


def _tb(slide, l, t, w, h, text, size=16, bold=False, color=WHITE, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    _run(p, text, size=size, bold=bold, color=color)
    return box


def _header(slide, title):
    _rect(slide, 0, 0, 13.333, 0.82, BLUE)
    _tb(slide, 0.35, 0.2, 12.5, 0.45, title, 22, True, WHITE)


def _table(slide, data, l, t, w, h):
    if not data:
        return
    rows, cols = len(data), len(data[0])
    tbl = slide.shapes.add_table(rows, cols, Inches(l), Inches(t), Inches(w), Inches(h)).table
    for j in range(cols):
        tbl.columns[j].width = Inches(w / cols)
    for i, row in enumerate(data):
        for j, val in enumerate(row):
            cell = tbl.cell(i, j)
            cell.text = ""
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER if j else PP_ALIGN.LEFT
            head = i == 0
            tot = str(row[0]).upper() == "TOTAL"
            _run(p, val, size=11 if not head else 12, bold=head or tot, color=WHITE if head or tot else SLATE)
            cell.fill.solid()
            cell.fill.fore_color.rgb = BLUE if head else (TEAL if tot else (LIGHT if i % 2 else WHITE))


def build_meeting_pptx(pack):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    isp = pack.get("isp", "ISP")
    rng = pack.get("range", "")

    s = prs.slides.add_slide(blank)
    _rect(s, 0, 0, 13.333, 7.5, NAVY)
    _rect(s, 0, 0, 0.16, 7.5, GOLD)
    _tb(s, 0.6, 1.5, 12, 0.35, "XTRNATE  •  NOC GOVERNANCE REVIEW", 14, True, GOLD)
    _tb(s, 0.6, 2.05, 12, 1.0, f"{isp}  CONCLUSION REPORT", 34, True, WHITE)
    _tb(s, 0.6, 3.2, 12, 0.4, rng, 18, False, MUTED)
    _tb(s, 0.6, 6.4, 12, 0.3, "Confidential  |  Partner review meeting", 12, False, MUTED)

    s = prs.slides.add_slide(blank)
    _header(s, f"{isp}  |  Snapshot")
    cards = pack.get("kpis", [])
    for i, (lab, val) in enumerate(cards[:5]):
        left = 0.35 + i * 2.58
        _rect(s, left, 1.25, 2.42, 1.55, [BLUE, TEAL, RGBColor(0x1D, 0x4E, 0xD8), RGBColor(0x7C, 0x3A, 0xED), RGBColor(0xB4, 0x53, 0x09)][i])
        _tb(s, left + 0.08, 1.38, 2.26, 0.3, lab, 11, True, MUTED, PP_ALIGN.CENTER)
        _tb(s, left + 0.08, 1.75, 2.26, 0.7, str(val), 26, True, WHITE, PP_ALIGN.CENTER)
    _tb(s, 0.4, 3.1, 12.5, 3.6, pack.get("snapshot_note", ""), 15, False, SLATE)

    s = prs.slides.add_slide(blank)
    _header(s, "What is working")
    _tb(s, 0.5, 1.2, 12.3, 5.5, pack.get("goods", "-"), 18, False, SLATE)

    s = prs.slides.add_slide(blank)
    _header(s, "Gaps — where focus is needed")
    _tb(s, 0.5, 1.2, 12.3, 5.5, pack.get("gaps", "-"), 18, False, SLATE)

    s = prs.slides.add_slide(blank)
    _header(s, "Outage classification (Last Remark)")
    _table(s, pack.get("class_table", [["Category", "Count", "%"]]), 0.35, 1.15, 12.6, 5.8)

    s = prs.slides.add_slide(blank)
    _header(s, "Remark tags — vendor / migration / feasibility")
    _table(s, pack.get("tag_table", [["Tag", "Count", "%"]]), 0.35, 1.15, 12.6, 5.8)

    s = prs.slides.add_slide(blank)
    _header(s, "State hotspots")
    _table(s, pack.get("state_table", [["State", "Count", "%"]]), 0.35, 1.15, 12.6, 5.8)

    s = prs.slides.add_slide(blank)
    _header(s, "Repeat sites (3 month / 6 month)")
    _table(s, pack.get("repeat_table", [["Site", "Period", "3M", "6M", "6M DT Hrs"]]), 0.3, 1.15, 12.7, 5.8)

    s = prs.slides.add_slide(blank)
    _header(s, "Top sites by downtime")
    _table(s, pack.get("site_table", [["Site", "Tickets", "DT Hrs", "Reasons"]]), 0.3, 1.15, 12.7, 5.8)

    s = prs.slides.add_slide(blank)
    _header(s, "Ageing / resolution bands")
    _table(s, pack.get("sla_table", [["Band", "Count", "%"]]), 0.35, 1.15, 12.6, 5.8)

    s = prs.slides.add_slide(blank)
    _header(s, "Open calls — risk now")
    _tb(s, 0.45, 1.15, 12.4, 1.0, pack.get("open_note", ""), 16, False, SLATE)
    _table(s, pack.get("open_table", [["Ticket", "Site", "Status", "State", "Hrs"]]), 0.3, 2.2, 12.7, 4.7)

    s = prs.slides.add_slide(blank)
    _header(s, "Vendor change / ISP change / Migration cases")
    _table(s, pack.get("change_table", [["Type", "Tickets", "Sites"]]), 0.35, 1.15, 12.6, 5.8)

    s = prs.slides.add_slide(blank)
    _header(s, "Where to focus next 30 days")
    _tb(s, 0.5, 1.2, 12.3, 5.5, pack.get("focus", "-"), 18, False, SLATE)

    s = prs.slides.add_slide(blank)
    _rect(s, 0, 0, 13.333, 7.5, NAVY)
    _rect(s, 0, 0, 0.16, 7.5, GOLD)
    _tb(s, 0.6, 2.3, 12, 0.5, "Action requested from partner", 26, True, WHITE)
    _tb(s, 0.6, 3.1, 12, 2.4, pack.get("ask", "-"), 18, False, MUTED)
    _tb(s, 0.6, 6.4, 12, 0.3, "XTRNATE Project  •  Hughes NOC", 13, False, GOLD)

    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()
