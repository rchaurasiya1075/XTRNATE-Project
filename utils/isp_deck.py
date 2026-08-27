"""Professional ISP review PowerPoint."""
from io import BytesIO
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsmap
from lxml import etree

NAVY = RGBColor(0x0F, 0x4C, 0x81)
NAVY2 = RGBColor(0x0B, 0x1F, 0x3A)
GOLD = RGBColor(0xD4, 0xA3, 0x37)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SLATE = RGBColor(0x1E, 0x29, 0x3B)
TEAL = RGBColor(0x0D, 0x94, 0x88)
LIGHT = RGBColor(0xF1, 0xF5, 0xF9)
MUTED = RGBColor(0x64, 0x74, 0x8B)


def _set_run(run, text, size=14, bold=False, color=WHITE, font="Calibri"):
    run.text = str(text)
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font


def _box(slide, l, t, w, h, fill, text="", size=14, bold=False, color=WHITE, align=PP_ALIGN.LEFT):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.fill.background()
    tf = sh.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.alignment = align
    _set_run(p.add_run() if p.runs else p.runs[0] if False else type("R", (), {})(), "")
    # simpler:
    p.clear() if hasattr(p, "clear") else None
    run = p.add_run()
    _set_run(run, text, size=size, bold=bold, color=color)
    return sh


def _rect(slide, l, t, w, h, fill):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.fill.background()
    return sh


def _tb(slide, l, t, w, h, text, size=16, bold=False, color=WHITE, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    _set_run(run, text, size=size, bold=bold, color=color)
    return box


def _add_table(slide, data, l, t, w, h, header=True):
    rows, cols = len(data), len(data[0]) if data else (0, 0)
    if rows == 0:
        return
    table = slide.shapes.add_table(rows, cols, Inches(l), Inches(t), Inches(w), Inches(h)).table
    for j in range(cols):
        table.columns[j].width = Inches(w / cols)
    for i, row in enumerate(data):
        for j, val in enumerate(row):
            cell = table.cell(i, j)
            cell.text = ""
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER if j else PP_ALIGN.LEFT
            run = p.add_run()
            is_head = header and i == 0
            is_tot = str(row[0]).upper() == "TOTAL"
            _set_run(
                run, val, size=11 if not is_head else 12,
                bold=is_head or is_tot,
                color=WHITE if is_head or is_tot else SLATE,
            )
            fill = NAVY if is_head else (TEAL if is_tot else (LIGHT if i % 2 else WHITE))
            cell.fill.solid()
            cell.fill.fore_color.rgb = fill
    return table


def build_isp_pptx(meta, cls_df, daily_df, state_df, site_df):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # 1 Cover
    s = prs.slides.add_slide(blank)
    _rect(s, 0, 0, 13.333, 7.5, NAVY2)
    _rect(s, 0, 0, 0.18, 7.5, GOLD)
    _tb(s, 0.7, 1.6, 12, 0.4, "XTRNATE  |  NOC COMMAND CENTER", 16, True, GOLD)
    _tb(s, 0.7, 2.1, 12, 1.1, f"{meta['isp']}  PERFORMANCE REVIEW", 36, True, WHITE)
    _tb(
        s, 0.7, 3.3, 12, 0.5,
        f"{meta['from']}   →   {meta['to']}     •     Date basis: {meta['date_on']}",
        18, False, RGBColor(0xCB, 0xD5, 0xE1),
    )
    _tb(s, 0.7, 6.4, 12, 0.4, "Confidential • For ISP review meeting", 12, False, MUTED)

    # 2 KPI
    s = prs.slides.add_slide(blank)
    _rect(s, 0, 0, 13.333, 0.85, NAVY)
    _tb(s, 0.4, 0.22, 12, 0.45, f"{meta['isp']}  |  Snapshot", 24, True, WHITE)
    cards = [
        ("TICKETS", meta["tickets"]),
        ("DOWNTIME HRS", meta["dt_hrs"]),
        ("AVG RESOLVE HRS", meta["avg_hrs"]),
        ("UNIQUE SITES", meta["sites"]),
        ("OPEN NOW", meta["open"]),
    ]
    colors = [NAVY, TEAL, RGBColor(0x1D, 0x4E, 0xD8), RGBColor(0x7C, 0x3A, 0xED), RGBColor(0xB4, 0x53, 0x09)]
    for i, ((lab, val), col) in enumerate(zip(cards, colors)):
        left = 0.35 + i * 2.58
        _rect(s, left, 1.4, 2.4, 1.7, col)
        _tb(s, left + 0.1, 1.55, 2.2, 0.35, lab, 11, True, RGBColor(0xE2, 0xE8, 0xF0), PP_ALIGN.CENTER)
        _tb(s, left + 0.1, 1.95, 2.2, 0.8, str(val), 28, True, WHITE, PP_ALIGN.CENTER)
    _tb(s, 0.4, 3.5, 12, 0.4, "Use this slide as the opening scorecard with the partner.", 14, False, SLATE)

    # 3 Classification
    s = prs.slides.add_slide(blank)
    _rect(s, 0, 0, 13.333, 0.85, NAVY)
    _tb(s, 0.4, 0.22, 12, 0.45, "Classification of Outage  •  from Last Remark", 22, True, WHITE)
    data = [["Outage Category", "Count", "%"]]
    if cls_df is not None and not cls_df.empty:
        for _, r in cls_df.head(16).iterrows():
            data.append([str(r.iloc[0])[:48], str(int(r.iloc[1]) if str(r.iloc[1]).replace('.','',1).isdigit() or isinstance(r.iloc[1], (int, float)) else r.iloc[1]), str(r.iloc[2])])
    _add_table(s, data, 0.4, 1.15, 12.5, min(5.8, 0.38 * len(data) + 0.3), header=True)

    # 4 State
    s = prs.slides.add_slide(blank)
    _rect(s, 0, 0, 13.333, 0.85, NAVY)
    _tb(s, 0.4, 0.22, 12, 0.45, "State-wise Outage Count", 22, True, WHITE)
    data = [["State", "Count", "%"]]
    if state_df is not None and not state_df.empty:
        for _, r in state_df.head(16).iterrows():
            data.append([str(r.iloc[0]), str(r.iloc[1]), str(r.iloc[2])])
    _add_table(s, data, 0.4, 1.15, 12.5, min(5.8, 0.38 * max(len(data), 2) + 0.3), header=True)

    # 5 Daily
    s = prs.slides.add_slide(blank)
    _rect(s, 0, 0, 13.333, 0.85, NAVY)
    _tb(s, 0.4, 0.22, 12, 0.45, "Daily Ticket Count", 22, True, WHITE)
    data = [["Date", "Count"]]
    if daily_df is not None and not daily_df.empty:
        tmp = daily_df.copy()
        tmp["Date"] = tmp["Date"].astype(str).str.slice(0, 10)
        for _, r in tmp.iterrows():
            data.append([str(r["Date"]), str(int(r["Count"]))])
    _add_table(s, data[:18], 0.4, 1.15, 8.5, min(5.8, 0.35 * min(len(data), 18) + 0.3), header=True)

    # 6 Sites
    s = prs.slides.add_slide(blank)
    _rect(s, 0, 0, 13.333, 0.85, NAVY)
    _tb(s, 0.4, 0.22, 12, 0.45, "Site Codes in Period  (Top 14 by ticket count)", 22, True, WHITE)
    data = [["Site Code", "Tickets", "State", "Category"]]
    if site_df is not None and not site_df.empty:
        cols = list(site_df.columns)
        for _, r in site_df.head(14).iterrows():
            sc = str(r.get("site_code", r.iloc[0]))
            tk = str(r.get("tickets", r.iloc[1] if len(r) > 1 else ""))
            stt = str(r.get("state", ""))[:28]
            cat = str(r.get("category", ""))[:42]
            data.append([sc, tk, stt, cat])
    _add_table(s, data, 0.35, 1.15, 12.6, min(5.8, 0.36 * len(data) + 0.25), header=True)

    # 7 Close
    s = prs.slides.add_slide(blank)
    _rect(s, 0, 0, 13.333, 7.5, NAVY2)
    _rect(s, 0, 0, 0.18, 7.5, GOLD)
    _tb(s, 0.7, 2.2, 12, 0.6, "Discussion points", 28, True, WHITE)
    _tb(
        s, 0.7, 3.0, 12, 2.2,
        "• High-repeat sites — joint action plan\n"
        "• Ageing open calls — revised ETR on MARS\n"
        "• Fibre / backend share — prevention + last-mile focus\n"
        "• Next review with same date-range format",
        18, False, RGBColor(0xE2, 0xE8, 0xF0),
    )
    _tb(s, 0.7, 6.3, 12, 0.35, "XTRNATE Project  •  Hughes NOC", 13, False, GOLD)

    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()
