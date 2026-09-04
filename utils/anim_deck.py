"""Animated graph PPT — GIFs + native charts + fade-in. Open with Slideshow (F5)."""
from __future__ import annotations

import os
import tempfile
from io import BytesIO
import pandas as pd
from lxml import etree
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

GREEN = RGBColor(0x1B, 0x4D, 0x3E)
GREEN2 = RGBColor(0x15, 0x3D, 0x31)
CREAM = RGBColor(0xF3, 0xF5, 0xF4)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
INK = RGBColor(0x1F, 0x1F, 0x1F)
MUTED = RGBColor(0xD4, 0xD9, 0xD6)
GOLD = RGBColor(0xC4, 0xA3, 0x5A)

PNS = "http://schemas.openxmlformats.org/presentationml/2006/main"


def _run(p, text, size=14, bold=False, color=WHITE, font="Calibri"):
    r = p.add_run()
    r.text = str(text)
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = font
    return r


def _rect(slide, l, t, w, h, fill):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.fill.background()
    return sh


def _tb(slide, l, t, w, h, text, size=16, bold=False, color=WHITE, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    _run(p, text, size=size, bold=bold, color=color)
    return box


def _header(slide, title):
    _rect(slide, 0, 0, 13.333, 0.78, GREEN)
    _tb(slide, 0.4, 0.18, 12.5, 0.45, title, 22, True, WHITE)


def _series(df, name_col, val_col, n=12):
    if df is None or getattr(df, "empty", True):
        return [], []
    d = df.copy()
    cols = list(d.columns)
    if name_col not in d.columns:
        name_col = cols[0]
    if val_col not in d.columns:
        # first numeric
        val_col = None
        for c in cols[1:]:
            if pd.api.types.is_numeric_dtype(d[c]) or c.lower() in ("count", "tickets", "%"):
                val_col = c
                break
        if val_col is None and len(cols) > 1:
            val_col = cols[1]
    if val_col is None:
        return [], []
    names = [str(x)[:28] for x in d[name_col].head(n).tolist()]
    vals = pd.to_numeric(d[val_col].head(n), errors="coerce").fillna(0).tolist()
    return names, [float(v) for v in vals]


def _gif_bars(labels, values, title) -> bytes | None:
    if not labels or not values:
        return None
    try:
        import matplotlib
        matplotlib.use("Agg")
        import matplotlib.pyplot as plt
        from matplotlib.animation import FuncAnimation, PillowWriter
        import numpy as np
    except Exception:
        return None
    labels = labels[:10]
    values = values[:10]
    x = np.arange(len(labels))
    frames = 16
    fig, ax = plt.subplots(figsize=(12.4, 5.4), dpi=110)
    fig.patch.set_facecolor("#1B4D3E")

    def draw(i):
        ax.clear()
        ax.set_facecolor("#1B4D3E")
        frac = (i + 1) / frames
        h = [v * frac for v in values]
        ax.bar(x, h, color="#E8F0EC", edgecolor="#C4A35A", linewidth=0.6)
        ax.set_xticks(x)
        ax.set_xticklabels(labels, rotation=25, ha="right", color="white", fontsize=9)
        ax.tick_params(colors="white")
        ax.set_ylim(0, max(values) * 1.18 if max(values) else 1)
        ax.set_title(title, color="white", fontsize=14, pad=8, fontweight="bold")
        for spine in ax.spines.values():
            spine.set_color("#2D6A4F")
        ax.yaxis.grid(True, color="#2D6A4F", alpha=0.5)
        ax.set_axisbelow(True)

    anim = FuncAnimation(fig, draw, frames=frames, interval=90)
    fd, path = tempfile.mkstemp(suffix=".gif")
    os.close(fd)
    try:
        anim.save(path, writer=PillowWriter(fps=10))
        with open(path, "rb") as f:
            return f.read()
    except Exception:
        return None
    finally:
        plt.close(fig)
        try:
            os.remove(path)
        except Exception:
            pass


def _gif_line(labels, values, title) -> bytes | None:
    if not labels or not values:
        return None
    try:
        import matplotlib
        matplotlib.use("Agg")
        import matplotlib.pyplot as plt
        from matplotlib.animation import FuncAnimation, PillowWriter
        import numpy as np
    except Exception:
        return None
    if len(values) > 28:
        idx = np.linspace(0, len(values) - 1, 28).astype(int)
        labels = [labels[i] for i in idx]
        values = [values[i] for i in idx]
    frames = max(len(values), 8)
    fig, ax = plt.subplots(figsize=(12.4, 5.4), dpi=110)
    fig.patch.set_facecolor("#1B4D3E")
    xs = list(range(len(values)))

    def draw(i):
        ax.clear()
        ax.set_facecolor("#1B4D3E")
        k = min(i + 1, len(values))
        ax.fill_between(xs[:k], values[:k], color="#2D6A4F", alpha=0.55)
        ax.plot(xs[:k], values[:k], color="#E8F0EC", lw=2.6, marker="o", ms=4)
        step = max(1, len(labels) // 8)
        ax.set_xticks(xs[::step])
        ax.set_xticklabels([str(labels[j])[:10] for j in xs[::step]], rotation=25, ha="right", color="white", fontsize=8)
        ax.tick_params(colors="white")
        ax.set_xlim(0, max(len(values) - 1, 1))
        ax.set_ylim(0, max(values) * 1.2 if max(values) else 1)
        ax.set_title(title, color="white", fontsize=14, pad=8, fontweight="bold")
        for spine in ax.spines.values():
            spine.set_color("#2D6A4F")
        ax.yaxis.grid(True, color="#2D6A4F", alpha=0.5)
        ax.set_axisbelow(True)

    anim = FuncAnimation(fig, draw, frames=frames, interval=120)
    fd, path = tempfile.mkstemp(suffix=".gif")
    os.close(fd)
    try:
        anim.save(path, writer=PillowWriter(fps=8))
        with open(path, "rb") as f:
            return f.read()
    except Exception:
        return None
    finally:
        plt.close(fig)
        try:
            os.remove(path)
        except Exception:
            pass


def _add_gif(slide, data: bytes, l, t, w, h):
    if not data:
        return None
    return slide.shapes.add_picture(BytesIO(data), Inches(l), Inches(t), Inches(w), Inches(h))


def _add_chart(slide, kind, cats, vals, series_name, l, t, w, h):
    if not cats or not vals:
        return None
    cd = CategoryChartData()
    cd.categories = cats
    cd.add_series(series_name, vals)
    ctype = XL_CHART_TYPE.LINE if kind == "line" else XL_CHART_TYPE.COLUMN_CLUSTERED
    gf = slide.shapes.add_chart(ctype, Inches(l), Inches(t), Inches(w), Inches(h), cd)
    chart = gf.chart
    try:
        chart.has_legend = False
        plot = chart.plots[0]
        plot.has_data_labels = False
        s = chart.series[0]
        s.format.fill.solid()
        s.format.fill.fore_color.rgb = GREEN
        if kind == "line":
            s.format.line.color.rgb = GREEN
    except Exception:
        pass
    return gf


def _fade_transition(slide):
    sld = slide._element
    for old in list(sld):
        if old.tag == qn("p:transition"):
            sld.remove(old)
    tr = etree.SubElement(sld, qn("p:transition"))
    tr.set("spd", "med")
    etree.SubElement(tr, qn("p:fade"))


def _appear_shapes(slide, skip=2):
    """Sequential fade-in (after previous) so Slideshow F5 animates."""
    ids = []
    for sh in slide.shapes:
        try:
            ids.append(int(sh.shape_id))
        except Exception:
            continue
    anim = ids[skip:] if len(ids) > skip else ids
    if not anim:
        return
    sld = slide._element
    for old in list(sld):
        if old.tag == qn("p:timing"):
            sld.remove(old)

    def el(tag, **attrs):
        node = etree.Element(qn(tag))
        for k, v in attrs.items():
            node.set(k, str(v))
        return node

    timing = el("p:timing")
    tnLst = el("p:tnLst")
    timing.append(tnLst)
    par_root = el("p:par")
    tnLst.append(par_root)
    cTn1 = el("p:cTn", id="1", dur="indefinite", restart="never", nodeType="tmRoot")
    par_root.append(cTn1)
    child1 = el("p:childTnLst")
    cTn1.append(child1)
    seq = el("p:seq", concurrent="true", nextAc="seek")
    child1.append(seq)
    cTn2 = el("p:cTn", id="2", restart="whenNotActive", fill="hold", evtFilter="cancelBubble", nodeType="mainSeq")
    seq.append(cTn2)
    child2 = el("p:childTnLst")
    cTn2.append(child2)

    nid = 3
    for i, sid in enumerate(anim):
        par = el("p:par")
        child2.append(par)
        cTn = el("p:cTn", id=str(nid), fill="hold")
        nid += 1
        par.append(cTn)
        st = el("p:stCondLst")
        cTn.append(st)
        delay = "0" if i == 0 else "400"
        st.append(el("p:cond", delay=delay))
        inner = el("p:childTnLst")
        cTn.append(inner)
        par2 = el("p:par")
        inner.append(par2)
        cTnB = el("p:cTn", id=str(nid), fill="hold")
        nid += 1
        par2.append(cTnB)
        st2 = el("p:stCondLst")
        cTnB.append(st2)
        st2.append(el("p:cond", delay="0"))
        ch = el("p:childTnLst")
        cTnB.append(ch)
        anim_el = el("p:animEffect", transition="in", filter="fade")
        ch.append(anim_el)
        cBhv = el("p:cBhv")
        anim_el.append(cBhv)
        cTnC = el("p:cTn", id=str(nid), dur="500")
        nid += 1
        cBhv.append(cTnC)
        tgt = el("p:tgtEl")
        cBhv.append(tgt)
        spTgt = el("p:spTgt", spid=str(sid))
        tgt.append(spTgt)

    prev = el("p:prevCondLst")
    seq.append(prev)
    prev.append(el("p:cond", evt="onPrev", delay="0"))
    nxt = el("p:nextCondLst")
    seq.append(nxt)
    nxt.append(el("p:cond", evt="onNext", delay="0"))
    sld.append(timing)


def _pct(part, whole):
    try:
        if not whole:
            return 0
        return round(100.0 * float(part) / float(whole), 1)
    except Exception:
        return 0


def _insights_daily(names, vals):
    if not vals:
        return ["Is period mein daily ticket series nahi mili.", "Date column check karke dubara export karo."]
    total = sum(vals)
    avg = total / len(vals)
    mx = max(vals)
    mn = min(vals)
    peak = names[vals.index(mx)] if names else "—"
    low = names[vals.index(mn)] if names else "—"
    first = sum(vals[: max(1, len(vals) // 3)])
    last = sum(vals[-max(1, len(vals) // 3) :])
    trend = "up" if last > first * 1.08 else ("down" if last < first * 0.92 else "flat")
    lines = [
        f"Period tickets: {int(total):,}  •  daily avg {avg:.1f}.",
        f"Peak day {peak} ({int(mx)}). Quietest {low} ({int(mn)}).",
    ]
    if trend == "up":
        lines.append("Trend UP — last third days pe load badha. Capacity / vendor backlog check.")
    elif trend == "down":
        lines.append("Trend DOWN — last third days pe tickets kam. Stabilisation dikh rahi hai.")
    else:
        lines.append("Trend FLAT — volume steady. Repeat sites pe nazar rakho.")
    if mx >= avg * 2 and avg > 0:
        lines.append("Peak 2× average se upar — spike ka root cause (fibre / node) alag se note karo.")
    lines.append("Graph play: line left→right draw hoti hai = din-by-din load.")
    return lines[:5]


def _insights_bars(names, vals, kind="class"):
    if not vals or not names:
        return ["Is cut pe data nahi. Filter / sheet check karo."]
    total = sum(vals) or 1
    order = sorted(zip(names, vals), key=lambda x: -x[1])
    top_n, top_v = order[0]
    share = _pct(top_v, total)
    lines = [f"Total in this view: {int(total):,} tickets."]
    lines.append(f"No.1 = {top_n}  ({int(top_v)}, {share}%).")
    if len(order) > 1:
        n2, v2 = order[1]
        lines.append(f"No.2 = {n2}  ({int(v2)}, {_pct(v2, total)}%).")
    top3 = _pct(sum(v for _, v in order[:3]), total)
    if kind == "class":
        lines.append(f"Top 3 classes = {top3}% of mix. Inhi pe restoration SOP tight karo.")
        if "fibre" in top_n.lower() or "fiber" in top_n.lower():
            lines.append("Fibre lead hai — route / MC / last-mile partner ko meeting agenda.")
        elif "vendor" in top_n.lower():
            lines.append("Vendor change lead — cutover window + LC update track karo.")
        elif "power" in top_n.lower():
            lines.append("Power lead — node UPS / feeder partner escalate.")
        else:
            lines.append("Lead class ko owner + TAT ke saath close karo.")
    elif kind == "state":
        lines.append(f"Top 3 states = {top3}% tickets. Field team yahin concentrate.")
        lines.append("High-count state mein repeat sites alag se nikaalo.")
    elif kind == "site":
        lines.append("Yeh chronic / hotspot sites hain — 3M/6M repeat ke saath padho.")
        lines.append("Top site pe last-mile + SIM backup status meeting mein confirm.")
    elif kind == "sla":
        lines.append("Green buckets = on-track. Long buckets = penalty risk.")
        lines.append(">24h share kam karna hi partner ask hona chahiye.")
    lines.append("Graph play: bars 0 se full height tak grow karte hain.")
    return lines[:6]


def _insights_exec(isp, rng, kpis, class_names, class_vals, daily_vals, state_names, state_vals):
    lines = [f"{isp}  •  {rng or 'selected period'}."]
    if kpis:
        bits = [f"{a}: {b}" for a, b in kpis[:4]]
        lines.append("  •  ".join(str(x) for x in bits))
    if class_names and class_vals:
        top = class_names[class_vals.index(max(class_vals))]
        lines.append(f"Dominant outage: {top} ({_pct(max(class_vals), sum(class_vals))}%).")
    if daily_vals:
        avg = sum(daily_vals) / len(daily_vals)
        lines.append(f"Daily load avg {avg:.1f} tickets (max {int(max(daily_vals))}).")
    if state_names and state_vals:
        st = state_names[state_vals.index(max(state_vals))]
        lines.append(f"Hottest state: {st}.")
    lines.append("Agle slides: animated graph + uski reading (visualization explain).")
    lines.append("Meeting use: F5 Slideshow, har graph ke right panel pe talking points.")
    return lines[:7]


def _panel(slide, l, t, w, h, title, lines):
    _rect(slide, l, t, w, h, GREEN2)
    _rect(slide, l, t, 0.08, h, GOLD)
    _tb(slide, l + 0.18, t + 0.12, w - 0.3, 0.4, title, 14, True, GOLD)
    body = "\n".join(f"•  {x}" for x in lines)
    _tb(slide, l + 0.18, t + 0.58, w - 0.32, h - 0.75, body, 13, False, WHITE)


def _viz_slide(prs, blank, title, gif, names, vals, kind, chart_kind="bar"):
    s = prs.slides.add_slide(blank)
    _header(s, title)
    if gif:
        _add_gif(s, gif, 0.25, 1.05, 8.5, 5.85)
    elif names and vals:
        _add_chart(s, chart_kind, names[:14], vals[:14], "Value", 0.25, 1.1, 8.5, 5.7)
    else:
        _tb(s, 0.5, 3.0, 8, 0.4, "No series for this graph.", 16, False, INK)
    explain = _insights_bars(names, vals, kind=kind) if kind != "daily" else _insights_daily(names, vals)
    _panel(s, 8.9, 1.1, 4.15, 5.8, "How to read / explain", explain)
    _fade_transition(s)
    return s


def build_animated_pptx(
    *,
    isp="ISP",
    rng="",
    kpis=None,
    class_df=None,
    daily_df=None,
    state_df=None,
    sla_df=None,
    site_df=None,
):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # Cover
    s = prs.slides.add_slide(blank)
    _rect(s, 0, 0, 13.333, 7.5, GREEN2)
    _rect(s, 0, 0, 0.16, 7.5, GOLD)
    _tb(s, 0.7, 1.7, 12, 0.35, "XTRNATE  •  ANIMATED PERFORMANCE DECK", 14, True, GOLD)
    _tb(s, 0.7, 2.2, 12, 1.0, f"{isp}", 36, True, WHITE)
    _tb(s, 0.7, 3.3, 12, 0.45, "Animated graphs  +  visualization explain  •  F5 Slideshow", 16, False, MUTED)
    _tb(s, 0.7, 3.85, 12, 0.4, rng, 16, False, WHITE)
    _tb(s, 0.7, 6.5, 12, 0.3, "Confidential  |  Partner review", 12, False, MUTED)
    _fade_transition(s)

    # KPI
    s = prs.slides.add_slide(blank)
    _header(s, f"{isp}  |  Snapshot")
    cards = kpis or []
    colors = [GREEN, RGBColor(0x2D, 0x6A, 0x4F), RGBColor(0x14, 0x53, 0x2D), RGBColor(0x3F, 0x5E, 0x56), RGBColor(0xB4, 0x53, 0x09)]
    for i, item in enumerate(cards[:5]):
        lab, val = item[0], item[1]
        left = 0.4 + i * 2.56
        _rect(s, left, 1.4, 2.4, 1.7, colors[i % len(colors)])
        _tb(s, left + 0.08, 1.52, 2.24, 0.35, str(lab), 11, True, MUTED, PP_ALIGN.CENTER)
        _tb(s, left + 0.08, 1.95, 2.24, 0.8, str(val), 26, True, WHITE, PP_ALIGN.CENTER)
    _tb(s, 0.5, 3.45, 12.4, 0.4, "F5 — cards fade. Agli slides pe graph + right side pe explain / report.", 14, False, INK)
    _fade_transition(s)
    _appear_shapes(s, skip=1)

    d_names, d_vals = [], []
    if daily_df is not None and not daily_df.empty:
        cols = list(daily_df.columns)
        date_c = cols[0]
        val_c = cols[1] if len(cols) > 1 else cols[0]
        for c in cols:
            cl = str(c).lower()
            if "date" in cl or "day" in cl:
                date_c = c
            if cl in ("count", "tickets", "ticket"):
                val_c = c
        tmp = daily_df.copy()
        tmp["_d"] = tmp[date_c].astype(str).str.slice(0, 10)
        tmp["_v"] = pd.to_numeric(tmp[val_c], errors="coerce").fillna(0)
        d_names = tmp["_d"].tolist()
        d_vals = tmp["_v"].astype(float).tolist()
    c_names, c_vals = _series(
        class_df,
        class_df.columns[0] if class_df is not None and len(class_df.columns) else "c",
        "Count",
    )
    st_names, st_vals = _series(
        state_df,
        state_df.columns[0] if state_df is not None and len(getattr(state_df, "columns", [])) else "State",
        "Count",
    )

    # Rewrite exec panel with full data — extra summary slide
    s = prs.slides.add_slide(blank)
    _header(s, f"{isp}  |  Executive visualization report")
    _panel(s, 0.35, 1.1, 12.6, 5.9, "What this deck shows", _insights_exec(
        isp, rng, cards, c_names, c_vals, d_vals, st_names, st_vals
    ))
    _fade_transition(s)

    gif = _gif_line(d_names, d_vals, "Daily ticket count")
    _viz_slide(prs, blank, "Daily tickets  •  animated line", gif, d_names, d_vals, "daily", "line")

    gif = _gif_bars(c_names, c_vals, "Tickets by outage class")
    _viz_slide(prs, blank, "Outage classification  •  bars grow", gif, c_names, c_vals, "class", "bar")

    gif = _gif_bars(st_names, st_vals, "Tickets by state")
    _viz_slide(prs, blank, "State-wise outage  •  bars grow", gif, st_names, st_vals, "state", "bar")

    # SLA
    if sla_df is not None and not sla_df.empty:
        names, vals = [], []
        labels0 = list(sla_df.iloc[:, 0].astype(str))
        if any(x.upper() in ("GRAND TOTAL", "TOTAL") for x in labels0):
            row = sla_df.iloc[-1]
            for c in sla_df.columns[1:]:
                v = pd.to_numeric(row[c], errors="coerce")
                if pd.notna(v):
                    names.append(str(c)[:18])
                    vals.append(float(v))
        else:
            names, vals = _series(sla_df, sla_df.columns[0], sla_df.columns[1] if len(sla_df.columns) > 1 else sla_df.columns[0])
        gif = _gif_bars(names[:8], vals[:8], "SLA mix")
        _viz_slide(prs, blank, "SLA buckets  •  bars grow", gif, names, vals, "sla", "bar")

    # Native editable charts slide
    s = prs.slides.add_slide(blank)
    _header(s, "Side-by-side visualization  (editable charts)")
    _add_chart(s, "bar", c_names[:8], c_vals[:8], "Class", 0.3, 1.05, 6.3, 4.3)
    _add_chart(s, "line", d_names[-14:], d_vals[-14:], "Daily", 6.8, 1.05, 6.15, 4.3)
    _panel(
        s, 0.3, 5.5, 12.7, 1.7,
        "Read together",
        [
            "Left = mix (kahan problem hai). Right = time (kab spike aaya).",
            "Agar spike + fibre same week hon = last-mile incident, sirf volume nahi.",
            "PowerPoint mein chart pe right-click → Animate for extra build.",
        ],
    )
    _fade_transition(s)
    _appear_shapes(s, skip=1)

    # Sites
    if site_df is not None and not site_df.empty:
        sn, sv = _series(site_df, site_df.columns[0], site_df.columns[1] if len(site_df.columns) > 1 else site_df.columns[0], n=10)
        gif = _gif_bars(sn, sv, "Top sites by tickets")
        _viz_slide(prs, blank, "Top sites  •  bars grow", gif, sn, sv, "site", "bar")

    s = prs.slides.add_slide(blank)
    _rect(s, 0, 0, 13.333, 7.5, GREEN2)
    _tb(s, 0.7, 1.8, 12, 0.45, "Close / partner ask", 14, True, GOLD)
    _tb(s, 0.7, 2.3, 12, 0.6, f"{isp}  —  actions from this visualization", 26, True, WHITE)
    actions = []
    if c_names and c_vals:
        actions.append(f"1. Own the lead class: {c_names[c_vals.index(max(c_vals))]}.")
    if st_names and st_vals:
        actions.append(f"2. Field push in {st_names[st_vals.index(max(st_vals))]}.")
    if d_vals:
        actions.append("3. Spike days ka RCA + repeat site list next review tak.")
    actions.append("4. SLA >24h tickets ko named owner + daily follow-up.")
    _tb(s, 0.7, 3.2, 12, 2.4, "\n".join(actions or ["Data ke hisaab se next review pe actions lock karo."]), 18, False, WHITE)
    _tb(s, 0.7, 6.5, 12, 0.3, "F5 Slideshow  •  Confidential  |  XTRNATE NOC", 12, False, MUTED)
    _fade_transition(s)

    out = BytesIO()
    prs.save(out)
    return out.getvalue()
